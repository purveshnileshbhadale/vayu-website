"""
file_processor.py — VAYU Universal File Processor

Supported types:
  image   → describe, ocr, resize, convert, compress, crop
  pdf     → summarize, extract_text, extract_pages, to_word
  docx    → summarize, extract_text, reformat, translate_hint
  txt/md  → summarize, reformat, translate_hint, word_count
  csv     → analyze, filter, sort, convert, stats
  xlsx    → analyze, filter, convert, stats
  json    → validate, format, extract, convert
  code    → explain, review, fix, run, document
  audio   → transcribe, trim, convert, info
  video   → trim, extract_audio, extract_frame, info, compress
  zip     → list, extract
  pptx    → summarize, extract_text, to_pdf
  html    → extract_text, summarize, to_md, to_docx, to_pdf, format, analyze
  epub    → extract_text, summarize
  rtf/odt → extract_text, summarize
  build   → build_docx, build_pptx, build_pdf, build_html, build_md
"""

import os
import re
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime

import google.generativeai as genai


def _get_api_key() -> str:
    config_path = Path(__file__).resolve().parent.parent / "config" / "api_keys.json"
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)["gemini_api_key"]


def _gemini_client():
    genai.configure(api_key=_get_api_key())
    return genai.GenerativeModel("gemini-2.5-flash")


def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "svg", "ico"}
    video_exts = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}
    audio_exts = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "wma", "opus"}
    code_exts  = {"py", "js", "ts", "jsx", "tsx", "css", "java", "c",
                  "cpp", "cs", "go", "rs", "rb", "php", "swift", "kt", "sh",
                  "bash", "ps1", "lua", "r", "m", "sql", "yaml", "toml"}
    html_exts  = {"html", "htm", "xhtml"}
    archive_exts = {"zip", "rar", "tar", "gz", "7z", "bz2", "xz"}
    doc_exts   = {"rtf", "odt", "ott"}

    if ext in image_exts:  return "image"
    if ext in video_exts:  return "video"
    if ext in audio_exts:  return "audio"
    if ext in code_exts:   return "code"
    if ext in html_exts:   return "html"
    if ext in archive_exts: return "archive"
    if ext in doc_exts:    return "richtext"
    if ext == "pdf":       return "pdf"
    if ext in ("docx", "doc"): return "docx"
    if ext in ("txt", "md", "rst", "log"): return "text"
    if ext in ("csv", "tsv"): return "csv"
    if ext in ("xlsx", "xls", "ods"): return "excel"
    if ext == "json":      return "json"
    if ext == "xml":       return "xml"
    if ext in ("pptx", "ppt"): return "pptx"
    if ext == "epub":      return "epub"
    return "unknown"


def _file_size_str(path: Path) -> str:
    size = path.stat().st_size
    if size < 1024:        return f"{size} B"
    if size < 1024**2:     return f"{size/1024:.1f} KB"
    if size < 1024**3:     return f"{size/1024**2:.1f} MB"
    return f"{size/1024**3:.1f} GB"

def _output_path(src: Path, suffix: str, new_ext: str = None) -> Path:
    ext  = new_ext or src.suffix
    name = f"{src.stem}_{suffix}{ext}"
    return src.parent / name

def _process_image(path: Path, action: str, params: dict, speak=None) -> str:
    try:
        from PIL import Image
    except ImportError:
        return "Pillow is not installed. Run: pip install Pillow"

    action = action or "describe"

    if action in ("describe", "ocr", "analyze", "read", "extract_text"):
        try:
            model  = _gemini_client()
            img    = Image.open(path)
            prompt = {
                "describe": "Describe this image in detail.",
                "ocr":      "Extract all text visible in this image. Return only the text, formatted clearly.",
                "analyze":  "Analyze this image thoroughly: objects, colors, composition, any text, context.",
                "read":     "Read all text in this image, preserving structure and formatting.",
                "extract_text": "Extract all text from this image.",
            }.get(action, "Describe this image.")

            if params.get("instruction"):
                prompt = params["instruction"]

            response = model.generate_content([prompt, img])
            result   = response.text.strip()

            if len(result) > 500 and params.get("save", True):
                out = _output_path(path, "result", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result[:300]}...\n\nFull result saved to: {out}"
            return result
        except Exception as e:
            return f"AI image analysis failed: {e}"

    if action == "resize":
        width  = int(params.get("width",  0))
        height = int(params.get("height", 0))
        scale  = float(params.get("scale", 0))
        try:
            img = Image.open(path)
            w, h = img.size
            if scale:
                new_size = (int(w * scale), int(h * scale))
            elif width and height:
                new_size = (width, height)
            elif width:
                new_size = (width, int(h * width / w))
            elif height:
                new_size = (int(w * height / h), height)
            else:
                return "Please specify width, height, or scale."
            out = _output_path(path, f"resized_{new_size[0]}x{new_size[1]}")
            img.resize(new_size, Image.LANCZOS).save(out)
            return f"Resized from {w}x{h} to {new_size[0]}x{new_size[1]}. Saved: {out.name}"
        except Exception as e:
            return f"Resize failed: {e}"

    if action == "convert":
        fmt = params.get("format", "png").lower().strip(".")
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG",
                   "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}
        pil_fmt = fmt_map.get(fmt, fmt.upper())
        try:
            img = Image.open(path).convert("RGB") if fmt == "jpg" else Image.open(path)
            out = _output_path(path, "converted", f".{fmt}")
            img.save(out, pil_fmt)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "compress":
        quality = int(params.get("quality", 70))
        try:
            img = Image.open(path).convert("RGB")
            out = _output_path(path, f"compressed_q{quality}", ".jpg")
            img.save(out, "JPEG", quality=quality, optimize=True)
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compressed: {before} → {after}. Saved: {out.name}"
        except Exception as e:
            return f"Compress failed: {e}"

    if action == "info":
        try:
            img = Image.open(path)
            return (f"Image info: {img.format}, {img.size[0]}x{img.size[1]}px, "
                    f"mode: {img.mode}, size: {_file_size_str(path)}")
        except Exception as e:
            return f"Info failed: {e}"

    return _process_image(path, "describe", {"instruction": f"{action}: {params}"})

def _process_pdf(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"

    def _extract_pdf_text(max_chars=50000) -> str:
        text = ""
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
        except ImportError:
            try:
                import PyPDF2
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            except ImportError:
                return ""
        return text[:max_chars]

    if action in ("summarize", "extract_text", "translate_hint", "analyze", "reformat"):
        text = _extract_pdf_text()
        if not text.strip():
            return "Could not extract text from PDF (may be scanned/image-based)."

        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(text, encoding="utf-8")
            return f"Text extracted ({len(text)} chars). Saved: {out.name}"

        prompt_map = {
            "summarize":      f"Summarize this PDF document concisely:\n\n{text}",
            "analyze":        f"Analyze this document thoroughly:\n\n{text}",
            "translate_hint": f"What language is this document in and what does it say? Summarize:\n\n{text}",
            "reformat":       f"Reformat this text cleanly with proper structure:\n\n{text}",
        }
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt_map.get(action, f"Analyze:\n\n{text}"))
            result   = response.text.strip()
            if len(result) > 600 and params.get("save", True):
                out = _output_path(path, action, ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result[:400]}...\n\nFull result saved: {out.name}"
            return result
        except Exception as e:
            return f"AI analysis failed: {e}"

    if action == "info":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
            return f"PDF: {pages} pages, size: {_file_size_str(path)}"
        except Exception:
            return f"PDF size: {_file_size_str(path)}"

    if action == "to_word":
        text = _extract_pdf_text()
        if not text:
            return "Could not extract text to convert."
        try:
            from docx import Document
            doc  = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "converted", ".docx")
            doc.save(out)
            return f"Converted to Word document. Saved: {out.name}"
        except ImportError:
            return "python-docx not installed. Run: pip install python-docx"

    return f"Unknown PDF action: '{action}'. Try: summarize, extract_text, info, to_word"

def _process_text_doc(path: Path, file_type: str, action: str,
                       params: dict, speak=None) -> str:
    action = action or "summarize"

    def _read_content() -> str:
        if file_type == "docx":
            try:
                from docx import Document
                doc  = Document(path)
                return "\n".join(p.text for p in doc.paragraphs)
            except ImportError:
                return "python-docx not installed."
            except Exception as e:
                return f"Read failed: {e}"
        else:
            return path.read_text(encoding="utf-8", errors="ignore")

    content = _read_content()
    if not content.strip():
        return "File appears to be empty."

    if action == "word_count":
        words = len(content.split())
        chars = len(content)
        lines = content.count("\n")
        return f"Word count: {words} words, {chars} characters, {lines} lines."

    if action == "extract_text":
        if file_type != "txt":
            out = _output_path(path, "extracted", ".txt")
            out.write_text(content, encoding="utf-8")
            return f"Text extracted. Saved: {out.name}"
        return content[:2000]

    instruction = params.get("instruction", "")
    prompt_map  = {
        "summarize":  f"Summarize this document concisely:\n\n{content[:40000]}",
        "analyze":    f"Analyze this document:\n\n{content[:40000]}",
        "reformat":   f"Reformat this text with clean structure, proper headings and paragraphs:\n\n{content[:40000]}",
        "fix":        f"Fix grammar, spelling and style issues in this text:\n\n{content[:40000]}",
        "translate_hint": f"What language is this and what does it say? Summarize:\n\n{content[:10000]}",
        "to_bullet":  f"Convert this text into a clear bullet-point summary:\n\n{content[:40000]}",
        "custom":     f"{instruction}\n\n{content[:40000]}",
    }

    if action not in prompt_map:

        action  = "custom"
        instruction = action

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt_map[action])
        result   = response.text.strip()
        if len(result) > 600 and params.get("save", True):
            out = _output_path(path, action, ".txt")
            out.write_text(result, encoding="utf-8")
            return f"{result[:400]}...\n\nFull result saved: {out.name}"
        return result
    except Exception as e:
        return f"AI processing failed: {e}"


def _process_html(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "extract_text"
    raw = path.read_text(encoding="utf-8", errors="ignore")

    if action == "extract_text":
        text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        out = _output_path(path, "extracted", ".txt")
        out.write_text(text, encoding="utf-8")
        return f"Text extracted ({len(text)} chars). Saved: {out.name}"

    if action == "format":
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(raw, 'html.parser')
            formatted = soup.prettify()
            out = _output_path(path, "formatted", ".html")
            out.write_text(formatted, encoding="utf-8")
            return f"HTML formatted. Saved: {out.name}"
        except ImportError:
            return "BeautifulSoup not installed. Run: pip install beautifulsoup4 html5lib"

    if action in ("summarize", "analyze"):
        text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()[:40000]
        prompt_map = {
            "summarize": f"Summarize this web page content concisely:\n\n{text}",
            "analyze":   f"Analyze this HTML page — structure, content, purpose, SEO:\n\nHTML:\n{raw[:5000]}\n\nText:\n{text}",
        }
        try:
            model = _gemini_client()
            response = model.generate_content(prompt_map.get(action, f"Analyze:\n\n{text}"))
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    if action == "to_md":
        try:
            from markdownify import markdownify
            cleaned = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
            cleaned = re.sub(r'<style[^>]*>.*?</style>', '', cleaned, flags=re.DOTALL)
            md = markdownify(cleaned, heading_style="ATX")
            out = _output_path(path, "converted", ".md")
            out.write_text(md, encoding="utf-8")
            return f"Converted to Markdown. Saved: {out.name}"
        except ImportError:
            return "markdownify not installed. Run: pip install markdownify"

    if action == "to_docx":
        text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split('\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "converted", ".docx")
            doc.save(out)
            return f"Converted to Word document. Saved: {out.name}"
        except ImportError:
            return "python-docx not installed."

    if action == "to_pdf":
        text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", "B", 16)
            pdf.cell(0, 10, path.stem, ln=True, align='C')
            pdf.ln(5)
            pdf.set_font("Arial", "", 11)
            for line in text.split('\n'):
                line = line.strip()
                if line:
                    safe = line.encode('latin-1', 'replace').decode('latin-1')
                    pdf.cell(0, 8, safe, ln=True)
            out = _output_path(path, "converted", ".pdf")
            pdf.output(str(out))
            return f"Converted to PDF. Saved: {out.name}"
        except ImportError:
            return "fpdf2 not installed. Run: pip install fpdf2"

    return f"Unknown HTML action: '{action}'. Try: extract_text, summarize, analyze, format, to_md, to_docx, to_pdf"


def _process_epub(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "extract_text"
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        return "ebooklib not installed. Run: pip install ebooklib beautifulsoup4"

    try:
        book = epub.read_epub(str(path))
    except Exception as e:
        return f"Could not read EPUB: {e}"

    chapters = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            html = item.get_body_content().decode('utf-8', errors='ignore')
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s+', ' ', text).strip()
            if text:
                chapters.append(text)

    full_text = '\n\n'.join(chapters)

    if action == "extract_text":
        out = _output_path(path, "extracted", ".txt")
        out.write_text(full_text, encoding="utf-8")
        return f"Text extracted ({len(full_text)} chars). Saved: {out.name}"

    if action in ("summarize", "analyze"):
        try:
            model = _gemini_client()
            prompt = f"{'Summarize' if action == 'summarize' else 'Analyze'} this ebook:\n\n{full_text[:40000]}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    return f"Unknown EPUB action: '{action}'. Try: extract_text, summarize"


def _process_richtext(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "extract_text"
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return "striprtf not installed. Run: pip install striprtf"

    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        text = rtf_to_text(raw)
    except Exception as e:
        return f"Could not read RTF: {e}"

    if action == "extract_text":
        out = _output_path(path, "extracted", ".txt")
        out.write_text(text, encoding="utf-8")
        return f"Text extracted ({len(text)} chars). Saved: {out.name}"

    if action in ("summarize", "analyze"):
        try:
            model = _gemini_client()
            prompt = f"{'Summarize' if action == 'summarize' else 'Analyze'} this document:\n\n{text[:40000]}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    return f"Unknown action: '{action}'. Try: extract_text, summarize"


def _build_document(action: str, params: dict, speak=None) -> str:
    content = params.get("content", "")
    title = params.get("title", "Document")
    output_path_str = params.get("output_path", "")
    instruction = params.get("instruction", "")

    if not content and instruction:
        try:
            model = _gemini_client()
            prompt = f"Generate document content. Title: {title}\n\nRequirement: {instruction}\n\nReturn only the document body content."
            response = model.generate_content(prompt)
            content = response.text.strip()
        except Exception as e:
            return f"AI content generation failed: {e}"

    if not content:
        return "No content or instruction provided to build document."

    if output_path_str:
        out = Path(output_path_str)
        out.parent.mkdir(parents=True, exist_ok=True)
    else:
        base = Path(".") if not params.get("save_path") else Path(params["save_path"])
        out = base / title

    if action == "build_docx":
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(title, 0)
            for para in content.split('\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = out.with_suffix('.docx')
            doc.save(out)
            return f"Document built ({len(content)} chars). Saved: {out}"
        except ImportError:
            return "python-docx not installed."

    if action == "build_pptx":
        try:
            from pptx import Presentation
            prs = Presentation()
            has_separator = '---' in content
            if has_separator:
                blocks = content.split('---')
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                for i, block in enumerate(blocks):
                    if i == 0:
                        continue
                    bullet_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_layout)
                    slide.shapes.title.text = block.strip().split('\n')[0][:50] if block.strip() else f"Slide {i}"
                    body = slide.shapes.placeholders[1]
                    tf = body.text_frame
                    for line in block.strip().split('\n')[1:]:
                        if line.strip():
                            p = tf.add_paragraph()
                            p.text = line.strip()
                slide_count = len(blocks) - 1
            else:
                bullet_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = title
                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()
                slide_count = 1
            out = out.with_suffix('.pptx')
            prs.save(out)
            return f"Presentation built ({slide_count} slides). Saved: {out}"
        except ImportError:
            return "python-pptx not installed."

    if action == "build_html":
        safe_title = title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        paragraphs = ''.join(f'<p>{p.strip()}</p>' for p in content.split('\n') if p.strip())
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{safe_title}</title>
    <style>
        body {{ font-family: 'Segoe UI', Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 2rem; line-height: 1.6; color: #1a1a1a; }}
        h1 {{ color: #2563eb; border-bottom: 2px solid #e5e7eb; padding-bottom: 0.5rem; }}
    </style>
</head>
<body>
    <h1>{safe_title}</h1>
    <div>{paragraphs}</div>
</body>
</html>"""
        out = out.with_suffix('.html')
        out.write_text(html, encoding="utf-8")
        return f"HTML page built. Saved: {out}"

    if action == "build_md":
        md = f"# {title}\n\n{content}"
        out = out.with_suffix('.md')
        out.write_text(md, encoding="utf-8")
        return f"Markdown file built. Saved: {out}"

    if action == "build_pdf":
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", "B", 16)
            safe_title = title.encode('latin-1', 'replace').decode('latin-1')
            pdf.cell(0, 10, safe_title, new_x="LMARGIN", new_y="NEXT", align='C')
            pdf.ln(5)
            pdf.set_font("Arial", "", 11)
            for line in content.split('\n'):
                line = line.strip()
                if line:
                    safe = line.encode('latin-1', 'replace').decode('latin-1')
                    pdf.set_x(pdf.l_margin)
                    pdf.multi_cell(0, 8, safe)
            out = out.with_suffix('.pdf')
            pdf.output(str(out))
            return f"PDF built. Saved: {out}"
        except ImportError:
            return "fpdf2 not installed. Run: pip install fpdf2"

    return f"Unknown build action: '{action}'. Try: build_docx, build_pptx, build_pdf, build_html, build_md"


def _process_data(path: Path, file_type: str, action: str,
                  params: dict, speak=None) -> str:
    try:
        import pandas as pd
    except ImportError:
        return "pandas not installed. Run: pip install pandas openpyxl"

    action = action or "analyze"

    try:
        if file_type == "csv":
            df = pd.read_csv(path, encoding="utf-8", errors="replace")
        else:
            df = pd.read_excel(path)
    except Exception as e:
        return f"Could not read file: {e}"

    if action == "info":
        return (f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                f"Columns: {', '.join(df.columns.tolist())}\n"
                f"Size: {_file_size_str(path)}")

    if action == "stats":
        try:
            desc = df.describe(include="all").to_string()
            return f"Statistics:\n{desc[:2000]}"
        except Exception as e:
            return f"Stats failed: {e}"

    if action == "analyze":
        preview = df.head(50).to_string()
        prompt  = (f"Analyze this dataset. Columns: {list(df.columns)}\n"
                   f"Rows: {len(df)}\nPreview:\n{preview}\n\n"
                   f"Give insights, patterns, and notable findings.")
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI analysis failed: {e}"

    if action in ("convert", "to_csv", "to_excel", "to_json"):
        fmt = {"to_csv": "csv", "to_excel": "xlsx", "to_json": "json",
               "convert": params.get("format", "csv")}.get(action, "csv")
        try:
            if fmt == "csv":
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False, encoding="utf-8")
            elif fmt == "xlsx":
                out = _output_path(path, "converted", ".xlsx")
                df.to_excel(out, index=False)
            elif fmt == "json":
                out = _output_path(path, "converted", ".json")
                df.to_json(out, orient="records", force_ascii=False, indent=2)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "filter":
        col       = params.get("column", "")
        value     = params.get("value", "")
        condition = params.get("condition", "equals")
        if not col or col not in df.columns:
            return f"Column '{col}' not found. Available: {', '.join(df.columns)}"
        try:
            if condition == "equals":     filtered = df[df[col] == value]
            elif condition == "contains": filtered = df[df[col].astype(str).str.contains(str(value), case=False)]
            elif condition == "gt":       filtered = df[df[col] > float(value)]
            elif condition == "lt":       filtered = df[df[col] < float(value)]
            else:                         filtered = df[df[col] == value]
            out = _output_path(path, "filtered", ".csv")
            filtered.to_csv(out, index=False)
            return f"Filtered: {len(filtered)} rows match. Saved: {out.name}"
        except Exception as e:
            return f"Filter failed: {e}"

    if action == "sort":
        col = params.get("column", df.columns[0])
        asc = params.get("ascending", True)
        try:
            sorted_df = df.sort_values(col, ascending=asc)
            out = _output_path(path, "sorted", path.suffix)
            sorted_df.to_csv(out, index=False)
            return f"Sorted by '{col}'. Saved: {out.name}"
        except Exception as e:
            return f"Sort failed: {e}"

    preview = df.head(30).to_string()
    try:
        model    = _gemini_client()
        response = model.generate_content(
            f"Task: {action}\nDataset ({len(df)} rows, cols: {list(df.columns)}):\n{preview}"
        )
        return response.text.strip()
    except Exception as e:
        return f"Processing failed: {e}"


def _process_json(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "analyze"
    try:
        content = path.read_text(encoding="utf-8")
        data    = json.loads(content)
    except Exception as e:
        return f"Invalid JSON: {e}"

    if action == "validate":
        return f"Valid JSON. Type: {type(data).__name__}, size: {_file_size_str(path)}"

    if action == "format":
        out = _output_path(path, "formatted", ".json")
        out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return f"Formatted JSON saved: {out.name}"

    if action in ("analyze", "summarize", "extract"):
        preview = json.dumps(data, indent=2, ensure_ascii=False)[:8000]
        prompt  = f"Task: {action} this JSON data:\n{preview}"
        if params.get("instruction"):
            prompt = f"{params['instruction']}\n\nJSON data:\n{preview}"
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    if action == "to_csv":
        try:
            import pandas as pd
            if isinstance(data, list):
                df  = pd.DataFrame(data)
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False)
                return f"Converted to CSV. Saved: {out.name}"
            return "JSON must be an array of objects to convert to CSV."
        except ImportError:
            return "pandas not installed."

    return _process_json(path, "analyze", {"instruction": action})

def _process_code(path: Path, action: str, params: dict, speak=None) -> str:
    action  = action or "explain"
    content = path.read_text(encoding="utf-8", errors="ignore")
    ext     = path.suffix.lstrip(".")

    if action == "run":
        if ext == "py":
            try:
                result = subprocess.run(
                    ["python", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                out = result.stdout or result.stderr
                return f"Output:\n{out[:2000]}" if out else "No output."
            except subprocess.TimeoutExpired:
                return "Execution timed out (30s)."
            except Exception as e:
                return f"Run failed: {e}"
        return f"Direct execution not supported for .{ext} files."

    if action == "info":
        lines = content.count("\n")
        words = len(content.split())
        return f"Code file: {lines} lines, {words} words, {_file_size_str(path)}"

    prompt_map = {
        "explain":   f"Explain this {ext} code clearly:\n\n```{ext}\n{content[:30000]}\n```",
        "review":    f"Review this {ext} code for bugs, issues, and improvements:\n\n```{ext}\n{content[:30000]}\n```",
        "fix":       f"Fix any bugs in this {ext} code and return the corrected version:\n\n```{ext}\n{content[:30000]}\n```",
        "optimize":  f"Optimize this {ext} code for performance and readability:\n\n```{ext}\n{content[:30000]}\n```",
        "document":  f"Add proper documentation/comments to this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
        "summarize": f"Summarize what this {ext} code does:\n\n```{ext}\n{content[:30000]}\n```",
        "test":      f"Write unit tests for this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
    }

    instruction = params.get("instruction", "")
    if action not in prompt_map:
        prompt = f"{action}\n\n```{ext}\n{content[:30000]}\n```"
        if instruction:
            prompt = f"{instruction}\n\n```{ext}\n{content[:30000]}\n```"
    else:
        prompt = prompt_map[action]

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt)
        result   = response.text.strip()

        if action in ("fix", "optimize", "document") and params.get("save", True):
            out = _output_path(path, action)
            code_match = re.search(r"```(?:\w+)?\n(.*?)```", result, re.DOTALL)
            code_to_save = code_match.group(1) if code_match else result
            out.write_text(code_to_save, encoding="utf-8")
            return f"{result[:400]}...\n\nSaved: {out.name}"
        return result
    except Exception as e:
        return f"AI processing failed: {e}"

def _process_audio(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "transcribe"

    if action == "info":
        try:
            from pydub import AudioSegment
            audio    = AudioSegment.from_file(path)
            duration = len(audio) / 1000
            mins, secs = divmod(int(duration), 60)
            return (f"Audio: {mins}m {secs}s, "
                    f"{audio.channels} ch, "
                    f"{audio.frame_rate}Hz, "
                    f"{_file_size_str(path)}")
        except ImportError:
            return f"Audio file: {_file_size_str(path)} (install pydub for more info)"
        except Exception as e:
            return f"Info failed: {e}"

    if action == "transcribe":
        try:
            model   = _gemini_client()
            content = path.read_bytes()
            mime    = {
                "mp3": "audio/mp3", "wav": "audio/wav",
                "ogg": "audio/ogg", "m4a": "audio/mp4",
                "aac": "audio/aac", "flac": "audio/flac",
            }.get(path.suffix.lstrip(".").lower(), "audio/mpeg")
            response = model.generate_content([
                "Transcribe all speech in this audio file accurately.",
                {"mime_type": mime, "data": content}
            ])
            result = response.text.strip()
            if params.get("save", True):
                out = _output_path(path, "transcript", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"Transcription saved: {out.name}\n\nPreview: {result[:300]}"
            return result
        except Exception as e:
            return f"Transcription failed: {e}"

    if action == "convert":
        fmt = params.get("format", "mp3").lstrip(".")
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            out   = _output_path(path, "converted", f".{fmt}")
            audio.export(out, format=fmt)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except ImportError:
            return "pydub not installed. Run: pip install pydub"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "trim":
        start = float(params.get("start", 0))
        end   = float(params.get("end",   0))
        try:
            from pydub import AudioSegment
            audio   = AudioSegment.from_file(path)
            end_ms  = int(end * 1000)   if end   else len(audio)
            trimmed = audio[int(start * 1000):end_ms]
            out     = _output_path(path, f"trim_{int(start)}s_{int(end)}s")
            trimmed.export(out, format=path.suffix.lstrip("."))
            return f"Trimmed audio ({int(start)}s–{int(end)}s). Saved: {out.name}"
        except ImportError:
            return "pydub not installed."
        except Exception as e:
            return f"Trim failed: {e}"

    return f"Unknown audio action: '{action}'. Try: transcribe, info, convert, trim"

def _process_video(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "info"


    def _ffmpeg_available() -> bool:
        try:
            subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=3)
            return True
        except Exception:
            return False

    if action == "info":
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json",
                 "-show_format", "-show_streams", str(path)],
                capture_output=True, text=True, timeout=10
            )
            data     = json.loads(result.stdout)
            fmt      = data.get("format", {})
            duration = float(fmt.get("duration", 0))
            mins, secs = divmod(int(duration), 60)
            size     = _file_size_str(path)
            streams  = data.get("streams", [])
            video_s  = next((s for s in streams if s["codec_type"] == "video"), {})
            w        = video_s.get("width", "?")
            h        = video_s.get("height", "?")
            fps      = video_s.get("r_frame_rate", "?")
            return f"Video: {mins}m {secs}s, {w}x{h}, {fps} fps, {size}"
        except Exception:
            return f"Video file: {_file_size_str(path)}"

    if action == "extract_audio":
        if not _ffmpeg_available():
            return "ffmpeg not found. Install ffmpeg to extract audio."
        out = _output_path(path, "audio", ".mp3")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a", str(out), "-y"],
                capture_output=True, timeout=300
            )
            return f"Audio extracted. Saved: {out.name}"
        except Exception as e:
            return f"Extract audio failed: {e}"

    if action == "trim":
        start = params.get("start", "00:00:00")
        end   = params.get("end",   "")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"trim", path.suffix)
        try:
            cmd = ["ffmpeg", "-i", str(path), "-ss", str(start)]
            if end:
                cmd += ["-to", str(end)]
            cmd += ["-c", "copy", str(out), "-y"]
            subprocess.run(cmd, capture_output=True, timeout=600)
            return f"Trimmed video saved: {out.name}"
        except Exception as e:
            return f"Trim failed: {e}"

    if action == "extract_frame":
        timestamp = params.get("timestamp", "00:00:01")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"frame_{timestamp.replace(':', '')}", ".jpg")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-ss", timestamp,
                 "-vframes", "1", str(out), "-y"],
                capture_output=True, timeout=30
            )
            return f"Frame extracted at {timestamp}. Saved: {out.name}"
        except Exception as e:
            return f"Extract frame failed: {e}"

    if action == "compress":
        crf = int(params.get("quality", 28))  
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"compressed_crf{crf}", ".mp4")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path),
                 "-c:v", "libx264", "-crf", str(crf),
                 "-preset", "medium", "-c:a", "copy",
                 str(out), "-y"],
                capture_output=True, timeout=1800
            )
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compressed: {before} → {after}. Saved: {out.name}"
        except Exception as e:
            return f"Compress failed: {e}"

    if action == "transcribe":
        if not _ffmpeg_available():
            return "ffmpeg not found. Needed for video transcription."
        tmp_audio = Path(tempfile.mktemp(suffix=".mp3"))
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a",
                 str(tmp_audio), "-y"],
                capture_output=True, timeout=300
            )
            result = _process_audio(tmp_audio, "transcribe", params, speak)
            return result
        except Exception as e:
            return f"Video transcription failed: {e}"
        finally:
            if tmp_audio.exists():
                tmp_audio.unlink()

    if action == "convert":
        fmt = params.get("format", "mp4").lstrip(".")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, "converted", f".{fmt}")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), str(out), "-y"],
                capture_output=True, timeout=1800
            )
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    return f"Unknown video action: '{action}'. Try: info, trim, extract_audio, extract_frame, compress, transcribe, convert"

def _process_archive(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "list"

    if action == "list":
        try:
            import zipfile, tarfile
            ext = path.suffix.lower()
            if ext == ".zip":
                with zipfile.ZipFile(path) as z:
                    names = z.namelist()
            elif ext in (".tar", ".gz", ".bz2", ".xz"):
                with tarfile.open(path) as t:
                    names = t.getnames()
            else:
                return f"Unsupported archive format: {ext}"
            preview = "\n".join(names[:30])
            suffix  = f"\n... and {len(names)-30} more" if len(names) > 30 else ""
            return f"Archive contains {len(names)} files:\n{preview}{suffix}"
        except Exception as e:
            return f"List failed: {e}"

    if action == "extract":
        dest = Path(params.get("destination", str(path.parent / path.stem)))
        dest.mkdir(parents=True, exist_ok=True)
        try:
            shutil.unpack_archive(path, dest)
            return f"Extracted to: {dest}"
        except Exception as e:
            return f"Extract failed: {e}"

    return f"Unknown archive action: '{action}'. Try: list, extract"

def _process_pptx(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"

    def _read_pptx_text() -> str:
        try:
            from pptx import Presentation
            prs  = Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                text.append(slide_text)
            return "\n".join(text)
        except ImportError:
            return "python-pptx not installed."

    if action in ("summarize", "extract_text", "analyze"):
        text = _read_pptx_text()
        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(text, encoding="utf-8")
            return f"Text extracted. Saved: {out.name}"
        try:
            model    = _gemini_client()
            prompt   = f"{'Summarize' if action == 'summarize' else 'Analyze'} this presentation:\n{text[:30000]}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    return f"Unknown PPTX action: '{action}'. Try: summarize, extract_text, analyze"

def file_processor(parameters: dict, player=None, speak=None) -> str:
    action = (parameters.get("action") or "").lower().strip()
    instruction = parameters.get("instruction", "")
    params = {**parameters, "instruction": instruction}

    if action.startswith("build_"):
        log_msg = f"[FileProcessor] BUILD | action={action}"
        print(log_msg)
        if player:
            player.write_log(log_msg)
        try:
            result = _build_document(action, params, speak)
            return result or "Done."
        except Exception as e:
            import traceback
            traceback.print_exc()
            return f"Build failed: {e}"

    file_path_str = parameters.get("file_path", "").strip()
    if not file_path_str:
        return "No file path provided."

    path = Path(file_path_str)
    if not path.exists():
        return f"File not found: {file_path_str}"
    if not path.is_file():
        return f"Path is not a file: {file_path_str}"

    file_type = _detect_type(path)

    log_msg = f"[FileProcessor] {file_type.upper()} | {path.name} | action={action or 'auto'}"
    print(log_msg)
    if player:
        player.write_log(log_msg)

    if file_type == "unknown":
        try:
            content = path.read_text(encoding="utf-8", errors="ignore")[:10000]
            model   = _gemini_client()
            prompt  = f"File: {path.name}\nContent preview:\n{content}\n\nTask: {action or instruction or 'Describe what this file contains and what can be done with it.'}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Unknown file type ({path.suffix}). Could not process: {e}"

    dispatch = {
        "image":    _process_image,
        "pdf":      _process_pdf,
        "docx":     lambda p, a, pm, s: _process_text_doc(p, "docx", a, pm, s),
        "text":     lambda p, a, pm, s: _process_text_doc(p, "text", a, pm, s),
        "html":     _process_html,
        "epub":     _process_epub,
        "richtext": _process_richtext,
        "csv":      lambda p, a, pm, s: _process_data(p, "csv",   a, pm, s),
        "excel":    lambda p, a, pm, s: _process_data(p, "excel", a, pm, s),
        "json":     _process_json,
        "xml":      lambda p, a, pm, s: _process_json(p, a, pm, s),
        "code":     _process_code,
        "audio":    _process_audio,
        "video":    _process_video,
        "archive":  _process_archive,
        "pptx":     _process_pptx,
    }

    handler = dispatch.get(file_type)
    if not handler:
        return f"Unsupported file type: {file_type}"

    try:
        result = handler(path, action, params, speak)
        return result or "Done."
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Processing failed: {e}"